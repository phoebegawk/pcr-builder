from io import BytesIO
from pathlib import Path
from datetime import datetime
import re
import copy

from fastapi import FastAPI, UploadFile, File, Form, HTTPException
from fastapi.responses import HTMLResponse, StreamingResponse
from fastapi.staticfiles import StaticFiles
from openpyxl import load_workbook
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image, UnidentifiedImageError
from pypdf import PdfReader
import uvicorn

from rep_data import REP_DATA, DROPDOWN_REPS

app = FastAPI(title="PCR Report Builder")

BASE_DIR = Path(__file__).resolve().parent
ASSETS_DIR = BASE_DIR / "assets"
TEMPLATES_DIR = BASE_DIR / "templates"
PPTX_TEMPLATE_PATH = TEMPLATES_DIR / "PCR - Template.pptx"

app.mount("/assets", StaticFiles(directory=str(ASSETS_DIR)), name="assets")


# -----------------------------
# Helpers
# -----------------------------
def clean_filename_part(value: str) -> str:
    value = value.strip()
    value = re.sub(r"[^A-Za-z0-9 _-]", "", value)
    value = re.sub(r"\s+", " ", value)
    return value or "Client"


def normalize_text(value: str) -> str:
    return re.sub(r"\s+", " ", str(value).strip())


def normalize_match_key(value: str) -> str:
    return re.sub(r"[^a-z0-9]+", "", str(value).lower())


def format_month_year(value) -> str:
    if isinstance(value, datetime):
        return value.strftime("%B %Y")

    if hasattr(value, "strftime"):
        return value.strftime("%B %Y")

    if isinstance(value, str) and value.strip():
        parsed = _parse_date_string(value.strip())
        if parsed:
            return parsed.strftime("%B %Y")

    raise ValueError("Couldn't format Month Year from the Excel ENDED date.")


def format_day_month_year(value) -> str:
    parsed = None

    if isinstance(value, datetime):
        parsed = value
    elif hasattr(value, "strftime"):
        parsed = value
    elif isinstance(value, str) and value.strip():
        parsed = _parse_date_string(value.strip())

    if not parsed:
        raise ValueError(f"Couldn't format date value: {value}")

    return f"{parsed.day} {parsed.strftime('%b')}, {parsed.year}"


def format_impressions(value) -> str:
    if value is None or str(value).strip() == "":
        return ""

    if isinstance(value, (int, float)):
        return f"{int(round(value)):,}"

    raw = str(value).strip().replace(",", "")
    try:
        return f"{int(float(raw)):,}"
    except ValueError:
        return str(value).strip()


def get_site_top_line(value) -> str:
    if value is None:
        return ""
    text = str(value).strip()
    if not text:
        return ""
    return text.splitlines()[0].strip()


def extract_month_year_from_excel(file_bytes: bytes) -> str:
    workbook = load_workbook(filename=BytesIO(file_bytes), data_only=True)

    for sheet in workbook.worksheets:
        ended_cell = None

        for row in sheet.iter_rows():
            for cell in row:
                if isinstance(cell.value, str) and cell.value.strip().upper() == "ENDED":
                    ended_cell = cell
                    break
            if ended_cell:
                break

        if not ended_cell:
            continue

        date_cell = sheet.cell(row=ended_cell.row + 1, column=ended_cell.column)
        return format_month_year(date_cell.value)

    raise ValueError("Couldn't find an 'ENDED' header in the uploaded Excel file.")


def extract_board_rows(file_bytes: bytes):
    workbook = load_workbook(filename=BytesIO(file_bytes), data_only=True)

    required_headers = {
        "SITE",
        "STARTED",
        "ENDED",
        "DAYS",
        "IMPRESSIONS",
    }

    for sheet in workbook.worksheets:
        header_row_idx = None
        header_map = {}

        for row in sheet.iter_rows():
            current_map = {}
            for cell in row:
                if isinstance(cell.value, str):
                    header_text = cell.value.strip().upper()
                    if header_text in required_headers:
                        current_map[header_text] = cell.column

            if required_headers.issubset(current_map.keys()):
                header_row_idx = row[0].row
                header_map = current_map
                break

        if not header_row_idx:
            continue

        rows = []
        current_row = header_row_idx + 1

        while current_row <= sheet.max_row:
            site_value = sheet.cell(current_row, header_map["SITE"]).value
            started_value = sheet.cell(current_row, header_map["STARTED"]).value
            ended_value = sheet.cell(current_row, header_map["ENDED"]).value
            days_value = sheet.cell(current_row, header_map["DAYS"]).value
            impressions_value = sheet.cell(current_row, header_map["IMPRESSIONS"]).value

            if all(
                value is None or str(value).strip() == ""
                for value in [site_value, started_value, ended_value, days_value, impressions_value]
            ):
                current_row += 1
                continue

            site_name = get_site_top_line(site_value)
            if not site_name:
                current_row += 1
                continue

            rows.append(
                {
                    "Site Name and Code": site_name,
                    "Start Date": format_day_month_year(started_value),
                    "End Date": format_day_month_year(ended_value),
                    "Run Time": f"{str(days_value).strip()} Days",
                    "Impressions": format_impressions(impressions_value),
                }
            )
            current_row += 1

        if rows:
            return rows

    raise ValueError(
        "Couldn't find a valid SITE / STARTED / ENDED / DAYS / IMPRESSIONS section in the uploaded Excel file."
    )


def _parse_date_string(value: str):
    formats = [
        "%d %B %Y",
        "%d %b %Y",
        "%d/%m/%Y",
        "%d-%m-%Y",
        "%Y-%m-%d",
    ]
    for fmt in formats:
        try:
            return datetime.strptime(value, fmt)
        except ValueError:
            pass
    return None


def extract_ado_contract_fields(pdf_bytes: bytes) -> dict:
    try:
        reader = PdfReader(BytesIO(pdf_bytes))
        text_parts = []
        for page in reader.pages:
            page_text = page.extract_text() or ""
            text_parts.append(page_text)
        full_text = "\n".join(text_parts)
    except Exception as exc:
        raise ValueError("Couldn't read the uploaded ADO PDF.") from exc

    if not full_text.strip():
        raise ValueError("The uploaded ADO PDF did not contain extractable text.")

    def extract_field(label: str) -> str:
        pattern = rf"{re.escape(label)}\s*:\s*(.+)"
        match = re.search(pattern, full_text, re.IGNORECASE)
        if not match:
            return ""
        return normalize_text(match.group(1))

    client_name = extract_field("Advertiser Name")
    sales_rep = extract_field("Account Manager")
    contract_number = extract_field("Contract Number")

    return {
        "client_name": client_name,
        "sales_rep": sales_rep,
        "contract_number": contract_number,
    }


def choose_contract_value(manual_value: str, ado_value: str, label: str) -> str:
    manual_value = normalize_text(manual_value or "")
    ado_value = normalize_text(ado_value or "")

    if manual_value:
        return manual_value
    if ado_value:
        return ado_value

    raise ValueError(f"{label} could not be found. Upload an ADO or enter it manually.")


def replace_text_on_slide(slide, replacements: dict):
    normalized_replacements = {
        normalize_text(key): value for key, value in replacements.items()
    }

    for shape in slide.shapes:
        if not hasattr(shape, "text_frame") or not shape.has_text_frame:
            continue

        current_text = normalize_text(shape.text)
        if current_text not in normalized_replacements:
            continue

        new_text = normalized_replacements[current_text]

        text_frame = shape.text_frame
        if not text_frame.paragraphs:
            text_frame.text = new_text
            continue

        first_paragraph = text_frame.paragraphs[0]
        if first_paragraph.runs:
            first_paragraph.runs[0].text = new_text
            for run in first_paragraph.runs[1:]:
                run.text = ""
        else:
            first_paragraph.text = new_text

        for paragraph in text_frame.paragraphs[1:]:
            for run in paragraph.runs:
                run.text = ""
            if not paragraph.runs:
                paragraph.text = ""


def duplicate_slide_safe(prs, slide_index: int):
    source = prs.slides[slide_index]
    blank_layout = prs.slide_layouts[6]
    new_slide = prs.slides.add_slide(blank_layout)

    for shape in source.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            image_stream = BytesIO(shape.image.blob)
            new_slide.shapes.add_picture(
                image_stream,
                shape.left,
                shape.top,
                width=shape.width,
                height=shape.height,
            )
        else:
            new_el = copy.deepcopy(shape.element)
            new_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")

    return new_slide


def move_slide_to_end(prs, slide_index: int):
    slide_id_list = prs.slides._sldIdLst
    slides = list(slide_id_list)
    slide = slides[slide_index]
    slide_id_list.remove(slide)
    slide_id_list.append(slide)


def find_board_placeholder_picture(slide, slide_width, slide_height):
    pictures = [shape for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE]
    if not pictures:
        return None

    slide_area = slide_width * slide_height
    filtered = [
        pic for pic in pictures
        if (pic.width * pic.height) < (slide_area * 0.80)
    ]

    candidates = filtered if filtered else pictures
    return max(candidates, key=lambda pic: pic.width * pic.height)


def delete_shape(shape):
    sp = shape.element
    sp.getparent().remove(sp)


def fit_image_within_bounds(image_bytes: bytes, bounds_width: int, bounds_height: int):
    try:
        with Image.open(BytesIO(image_bytes)) as img:
            img_width, img_height = img.size
    except UnidentifiedImageError as exc:
        raise ValueError("One of the uploaded image files could not be read.") from exc

    if img_width <= 0 or img_height <= 0:
        raise ValueError("One of the uploaded image files has invalid dimensions.")

    scale = min(bounds_width / img_width, bounds_height / img_height)
    fitted_width = int(img_width * scale)
    fitted_height = int(img_height * scale)

    return fitted_width, fitted_height


def replace_board_placeholder_image(slide, slide_width, slide_height, image_bytes: bytes):
    placeholder = find_board_placeholder_picture(slide, slide_width, slide_height)
    if not placeholder:
        return

    left = placeholder.left
    top = placeholder.top
    width = placeholder.width
    height = placeholder.height

    fitted_width, fitted_height = fit_image_within_bounds(image_bytes, width, height)

    new_left = left + int((width - fitted_width) / 2)
    new_top = top + int((height - fitted_height) / 2)

    delete_shape(placeholder)
    slide.shapes.add_picture(
        BytesIO(image_bytes),
        new_left,
        new_top,
        width=fitted_width,
        height=fitted_height,
    )


def collect_uploaded_images(board_images: list[UploadFile]):
    allowed_extensions = (".jpg", ".jpeg", ".png")
    collected = []

    for image in board_images:
        if not image.filename:
            continue

        if not image.filename.lower().endswith(allowed_extensions):
            raise ValueError(
                f"Invalid image file '{image.filename}'. Only JPG, JPEG and PNG are allowed."
            )

        collected.append(image)

    return collected


async def read_uploaded_images(board_images: list[UploadFile]):
    uploaded = []
    for image in collect_uploaded_images(board_images):
        image_bytes = await image.read()
        uploaded.append(
            {
                "filename": image.filename,
                "match_key": normalize_match_key(Path(image.filename).stem),
                "bytes": image_bytes,
            }
        )
    return uploaded


def find_matching_image_bytes(site_name: str, uploaded_images: list[dict]):
    site_key = normalize_match_key(site_name)
    if not site_key:
        return None

    for image in uploaded_images:
        if site_key in image["match_key"]:
            return image["bytes"]

    return None


def build_pcr_pptx(
    client_name: str,
    month_year: str,
    contract_number: str,
    sales_rep: str,
    board_rows: list,
    uploaded_images: list[dict],
) -> BytesIO:
    if not PPTX_TEMPLATE_PATH.exists():
        raise FileNotFoundError(f"Template not found at: {PPTX_TEMPLATE_PATH}")

    if sales_rep not in REP_DATA:
        raise ValueError("Selected sales rep wasn't recognised.")

    if not board_rows:
        raise ValueError("No board rows were found in the uploaded Excel file.")

    rep = REP_DATA[sales_rep]
    prs = Presentation(str(PPTX_TEMPLATE_PATH))

    if len(prs.slides) < 3:
        raise ValueError("The PPTX template must contain at least 3 slides: cover, board page, contact page.")

    cover_slide = prs.slides[0]
    board_template_index = 1
    contact_slide_index = 2

    replace_text_on_slide(
        cover_slide,
        {
            "Client Name": client_name,
            "Month Year": month_year,
            "Contract Number": contract_number,
        },
    )

    first_board_slide = prs.slides[board_template_index]
    replace_text_on_slide(first_board_slide, board_rows[0])

    matched_image = find_matching_image_bytes(board_rows[0]["Site Name and Code"], uploaded_images)
    if matched_image:
        replace_board_placeholder_image(
            first_board_slide,
            prs.slide_width,
            prs.slide_height,
            matched_image,
        )

    for row_data in board_rows[1:]:
        duplicated_board_slide = duplicate_slide_safe(prs, board_template_index)
        replace_text_on_slide(duplicated_board_slide, row_data)

        matched_image = find_matching_image_bytes(row_data["Site Name and Code"], uploaded_images)
        if matched_image:
            replace_board_placeholder_image(
                duplicated_board_slide,
                prs.slide_width,
                prs.slide_height,
                matched_image,
            )

    contact_slide = prs.slides[contact_slide_index]
    rep_contact_line = f'{rep["phone"]} | {rep["email"]}'

    replace_text_on_slide(
        contact_slide,
        {
            "Rep Name!": rep["display_name"],
            "Rep Number | Rep Email": rep_contact_line,
            "Rep Number  |  Rep Email": rep_contact_line,
            "Rep Number   |   Rep Email": rep_contact_line,
        },
    )

    move_slide_to_end(prs, contact_slide_index)

    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output


# -----------------------------
# Routes
# -----------------------------
@app.get("/", response_class=HTMLResponse)
async def home():
    options_html = "".join(
        f'<option value="{rep}"></option>' for rep in DROPDOWN_REPS
    )

    return f"""
    <html>
    <head>
        <title>PCR Builder</title>
        <link rel="icon" type="image/png" href="/assets/favicon.png">
        <link rel="stylesheet" href="/assets/styles.css">
        <meta name="viewport" content="width=device-width, initial-scale=1.0">
        <style>
            .client-name-input {{
                width: 100%;
                padding: 12px 14px;
                border-radius: 8px;
                border: 2px solid #D7DF23;
                background: #FFFFFF;
                color: #542D54;
                font-size: 16px;
                font-weight: 600;
            }}

            .field-note {{
                width: 100%;
                color: #542D54;
                font-size: 13px;
                font-weight: 600;
                opacity: 0.9;
                margin-top: 6px;
            }}

            .section-label {{
                width: 100%;
                color: #542D54;
                font-size: 24px;
                font-weight: 700;
                margin: 0 0 10px 0;
                line-height: 1.1;
            }}

            .section-block {{
                width: 100%;
            }}

            .section-block-spaced {{
                width: 100%;
                margin-top: 12px;
            }}

            .manual-section {{
                display: none;
                margin-top: 14px;
                padding-top: 10px;
                border-top: 2px dashed rgba(84, 45, 84, 0.25);
            }}

            .manual-section.visible {{
                display: block;
            }}

            .secondary-button {{
                background: #FFFFFF;
                color: #542D54;
                border: 3px solid #D7DF23;
                min-width: 220px;
            }}

            .gawk-button {{
                border: none;
                border-radius: 999px;
                padding: 14px 28px;
                font-size: 16px;
                font-weight: 700;
                cursor: pointer;
                transition: transform 0.15s ease, opacity 0.15s ease;
            }}

            .gawk-button:hover {{
                transform: translateY(-1px);
            }}

            .browse-btn {{
                background: #D7DF23;
                color: #542D54;
            }}

            .actions-row {{
                width: 100%;
                display: flex;
                justify-content: center;
                gap: 14px;
                margin: 8px auto 0 auto;
                flex-wrap: wrap;
            }}

            .build-button {{
                background: linear-gradient(90deg, #d9b3db 0%, #8b7f40 100%);
                color: #542D54;
                min-width: 170px;
            }}

            .reset-button {{
                display: block;
                margin: 24px auto 0 auto;
                background: #FFFFFF;
                color: #542D54;
                border: 3px solid #D7DF23;
                min-width: 140px;
            }}

            .upload-confirm-text {{
                color: #542D54;
                font-weight: 700;
            }}

            .page-subtext {{
                width: 100%;
                max-width: 1200px;
                margin: -18px auto 24px auto;
                text-align: center;
                color: #FFFFFF;
            }}

            .page-subtext-primary {{
                font-weight: 700;
                font-size: 16px;
                margin-bottom: 4px;
            }}

            .page-subtext-secondary {{
                font-weight: 600;
                font-size: 14px;
                opacity: 0.95;
            }}

            .status-message {{
                width: 100%;
                text-align: center;
                color: #542D54;
                font-weight: 700;
                font-size: 14px;
                min-height: 20px;
            }}

            .dropdown {{
                width: 100%;
                padding: 12px 14px;
                border-radius: 8px;
                border: 2px solid #D7DF23;
                background: #FFFFFF;
                color: #542D54;
                font-size: 16px;
                font-weight: 600;
            }}

            @media (max-width: 768px) {{
                body {{
                    padding: 20px;
                }}

                .spec-section {{
                    padding: 18px;
                }}

                .section-label {{
                    font-size: 20px;
                }}
            }}
        </style>
    </head>

    <body>
        <div class="header-container">
            <img src="/assets/Header-PCRBuilder.png" class="header-image" alt="PCR Builder">
        </div>

        <div class="page-subtext">
            <div class="page-subtext-primary">Build your PCR using the PCR Numbers Excel file.</div>
            <div class="page-subtext-secondary">Upload an ADO first, or use manual contract information entry.</div>
        </div>

        <form id="pcrForm" class="spec-section">
            <div class="section-inner">
                <div class="section-block">
                    <div class="section-label">Upload ADO PDF</div>
                    <div class="drop-area" id="adoDropArea">
                        <p>Drag & drop ADO PDF here!</p>
                        <button type="button" class="gawk-button browse-btn" id="adoBrowseBtn">Browse File</button>
                        <input
                            type="file"
                            id="adoFile"
                            class="file-input"
                            accept=".pdf"
                        />
                    </div>
                    <div class="field-note">Used to extract Client Name, Sales Representative, and Contract Number.</div>
                </div>

                <div class="actions-row" style="margin-top: 16px;">
                    <button type="button" id="manualToggleBtn" class="gawk-button secondary-button">Manual Contract Information Entry</button>
                </div>

                <div id="manualSection" class="manual-section">
                    <div class="section-block">
                        <div class="section-label">Client Name</div>
                        <input
                            type="text"
                            id="clientName"
                            name="client_name"
                            class="client-name-input"
                            placeholder=""
                            autocomplete="off"
                        />
                        <div class="field-note">Manual override for Client Name.</div>
                    </div>

                    <div class="section-block-spaced">
                        <div class="section-label">Sales Representative</div>
                        <input
                            type="text"
                            id="salesRep"
                            name="sales_rep"
                            class="dropdown"
                            list="salesRepOptions"
                            placeholder=""
                            autocomplete="off"
                        />
                        <datalist id="salesRepOptions">
                            {options_html}
                        </datalist>
                        <div class="field-note">Manual override for Sales Representative.</div>
                    </div>

                    <div class="section-block-spaced">
                        <div class="section-label">Contract Number</div>
                        <input
                            type="text"
                            id="contractNumber"
                            name="contract_number"
                            class="client-name-input"
                            placeholder=""
                            autocomplete="off"
                        />
                        <div class="field-note">Manual override for Contract Number.</div>
                    </div>
                </div>

                <div class="section-block-spaced">
                    <div class="section-label">Upload PCR Numbers File</div>
                    <div class="drop-area" id="dropArea">
                        <p>Drag & drop Excel file here!</p>
                        <button type="button" class="gawk-button browse-btn" id="browseBtn">Browse File</button>
                        <input
                            type="file"
                            id="excelFile"
                            class="file-input"
                            accept=".xlsx,.xlsm,.xltx,.xltm"
                            required
                        />
                    </div>
                </div>

                <div class="section-block-spaced">
                    <div class="section-label">Upload PoP's</div>
                    <div class="drop-area" id="imageDropArea">
                        <p>Drag & drop JPG and PNG here!</p>
                        <button type="button" class="gawk-button browse-btn" id="imageBrowseBtn">Browse Files</button>
                        <input
                            type="file"
                            id="imageFiles"
                            class="file-input"
                            accept=".jpg,.jpeg,.png"
                            multiple
                        />
                    </div>
                </div>

                <div class="upload-confirm hidden" id="adoUploadConfirm">
                    <div class="upload-confirm-text" id="adoUploadConfirmText">ADO ready.</div>
                </div>

                <div class="upload-confirm hidden" id="uploadConfirm">
                    <div class="upload-confirm-text" id="uploadConfirmText">File ready.</div>
                </div>

                <div class="upload-confirm hidden" id="imageUploadConfirm">
                    <div class="upload-confirm-text" id="imageUploadConfirmText">Images ready.</div>
                </div>

                <div class="status-message" id="statusMessage"></div>
            </div>
        </form>

        <div class="actions-row">
            <button id="buildBtn" class="gawk-button build-button">Build PCR</button>
        </div>

        <button id="resetBtn" class="gawk-button reset-button">Reset all</button>

        <script>
            const form = document.getElementById("pcrForm");
            const manualSection = document.getElementById("manualSection");
            const manualToggleBtn = document.getElementById("manualToggleBtn");

            const adoFileInput = document.getElementById("adoFile");
            const adoBrowseBtn = document.getElementById("adoBrowseBtn");
            const adoDropArea = document.getElementById("adoDropArea");
            const adoUploadConfirm = document.getElementById("adoUploadConfirm");
            const adoUploadConfirmText = document.getElementById("adoUploadConfirmText");

            const clientNameInput = document.getElementById("clientName");
            const salesRepSelect = document.getElementById("salesRep");
            const contractNumberInput = document.getElementById("contractNumber");

            const excelFileInput = document.getElementById("excelFile");
            const imageFilesInput = document.getElementById("imageFiles");
            const browseBtn = document.getElementById("browseBtn");
            const imageBrowseBtn = document.getElementById("imageBrowseBtn");
            const buildBtn = document.getElementById("buildBtn");
            const resetBtn = document.getElementById("resetBtn");
            const dropArea = document.getElementById("dropArea");
            const imageDropArea = document.getElementById("imageDropArea");
            const uploadConfirm = document.getElementById("uploadConfirm");
            const uploadConfirmText = document.getElementById("uploadConfirmText");
            const imageUploadConfirm = document.getElementById("imageUploadConfirm");
            const imageUploadConfirmText = document.getElementById("imageUploadConfirmText");
            const statusMessage = document.getElementById("statusMessage");

            manualToggleBtn.addEventListener("click", () => {{
                manualSection.classList.toggle("visible");
            }});

            adoBrowseBtn.addEventListener("click", () => adoFileInput.click());
            browseBtn.addEventListener("click", () => excelFileInput.click());
            imageBrowseBtn.addEventListener("click", () => imageFilesInput.click());

            adoFileInput.addEventListener("change", () => {{
                updateSelectedAdoFile();
            }});

            excelFileInput.addEventListener("change", () => {{
                updateSelectedExcelFile();
            }});

            imageFilesInput.addEventListener("change", () => {{
                updateSelectedImageFiles();
            }});

            ["dragenter", "dragover"].forEach(eventName => {{
                [adoDropArea, dropArea, imageDropArea].forEach(area => {{
                    area.addEventListener(eventName, (e) => {{
                        e.preventDefault();
                        e.stopPropagation();
                        area.classList.add("drag-over");
                    }});
                }});
            }});

            ["dragleave", "drop"].forEach(eventName => {{
                [adoDropArea, dropArea, imageDropArea].forEach(area => {{
                    area.addEventListener(eventName, (e) => {{
                        e.preventDefault();
                        e.stopPropagation();
                        area.classList.remove("drag-over");
                    }});
                }});
            }});

            adoDropArea.addEventListener("drop", (e) => {{
                const files = e.dataTransfer.files;
                if (!files || !files.length) return;

                const dt = new DataTransfer();
                dt.items.add(files[0]);
                adoFileInput.files = dt.files;
                updateSelectedAdoFile();
            }});

            dropArea.addEventListener("drop", (e) => {{
                const files = e.dataTransfer.files;
                if (!files || !files.length) return;

                const dt = new DataTransfer();
                dt.items.add(files[0]);
                excelFileInput.files = dt.files;
                updateSelectedExcelFile();
            }});

            imageDropArea.addEventListener("drop", (e) => {{
                const files = e.dataTransfer.files;
                if (!files || !files.length) return;

                const dt = new DataTransfer();
                for (let i = 0; i < files.length; i++) {{
                    dt.items.add(files[i]);
                }}
                imageFilesInput.files = dt.files;
                updateSelectedImageFiles();
            }});

            function updateSelectedAdoFile() {{
                if (adoFileInput.files && adoFileInput.files.length > 0) {{
                    adoUploadConfirm.classList.remove("hidden");
                    adoUploadConfirmText.textContent = `ADO ready: ${{adoFileInput.files[0].name}}`;
                    statusMessage.textContent = "";
                }} else {{
                    adoUploadConfirm.classList.add("hidden");
                }}
            }}

            function updateSelectedExcelFile() {{
                if (excelFileInput.files && excelFileInput.files.length > 0) {{
                    uploadConfirm.classList.remove("hidden");
                    uploadConfirmText.textContent = `File ready: ${{excelFileInput.files[0].name}}`;
                    statusMessage.textContent = "";
                }} else {{
                    uploadConfirm.classList.add("hidden");
                }}
            }}

            function updateSelectedImageFiles() {{
                if (imageFilesInput.files && imageFilesInput.files.length > 0) {{
                    imageUploadConfirm.classList.remove("hidden");
                    imageUploadConfirmText.textContent = `${{imageFilesInput.files.length}} image file(s) ready`;
                    statusMessage.textContent = "";
                }} else {{
                    imageUploadConfirm.classList.add("hidden");
                }}
            }}

            buildBtn.addEventListener("click", async () => {{
                const clientName = clientNameInput.value.trim();
                const salesRep = salesRepSelect.value.trim();
                const contractNumber = contractNumberInput.value.trim();
                const excelFile = excelFileInput.files[0];
                const adoFile = adoFileInput.files[0];

                if (!excelFile) {{
                    statusMessage.textContent = "Please upload a PCR Numbers Excel file.";
                    return;
                }}

                if (!adoFile && !clientName && !salesRep && !contractNumber) {{
                    statusMessage.textContent = "Upload an ADO PDF or use manual contract information entry.";
                    return;
                }}

                statusMessage.textContent = "Building PCR...";
                buildBtn.disabled = true;

                const formData = new FormData();
                formData.append("client_name", clientName);
                formData.append("sales_rep", salesRep);
                formData.append("contract_number", contractNumber);
                formData.append("excel_file", excelFile);

                if (adoFile) {{
                    formData.append("ado_file", adoFile);
                }}

                for (let i = 0; i < imageFilesInput.files.length; i++) {{
                    formData.append("board_images", imageFilesInput.files[i]);
                }}

                try {{
                    const response = await fetch("/build", {{
                        method: "POST",
                        body: formData
                    }});

                    if (!response.ok) {{
                        let errorText = "Something went wrong building the PCR.";
                        try {{
                            const errorJson = await response.json();
                            errorText = errorJson.detail || errorText;
                        }} catch {{
                            errorText = "Something went wrong building the PCR.";
                        }}
                        throw new Error(errorText);
                    }}

                    const blob = await response.blob();
                    const url = window.URL.createObjectURL(blob);
                    const anchor = document.createElement("a");
                    const disposition = response.headers.get("Content-Disposition") || "";
                    const match = disposition.match(/filename="([^"]+)"/);
                    anchor.href = url;
                    anchor.download = match ? match[1] : "PCR_Report.pptx";
                    document.body.appendChild(anchor);
                    anchor.click();
                    anchor.remove();
                    window.URL.revokeObjectURL(url);

                    statusMessage.textContent = "Done. Your PCR PPTX has been downloaded.";
                }} catch (error) {{
                    statusMessage.textContent = error.message || "Failed to build PCR.";
                }} finally {{
                    buildBtn.disabled = false;
                }}
            }});

            resetBtn.addEventListener("click", () => {{
                form.reset();
                adoFileInput.value = "";
                excelFileInput.value = "";
                imageFilesInput.value = "";
                adoUploadConfirm.classList.add("hidden");
                uploadConfirm.classList.add("hidden");
                imageUploadConfirm.classList.add("hidden");
                statusMessage.textContent = "";
            }});
        </script>
    </body>
    </html>
    """


@app.post("/build")
async def build_pcr(
    client_name: str = Form(""),
    sales_rep: str = Form(""),
    contract_number: str = Form(""),
    excel_file: UploadFile = File(...),
    ado_file: UploadFile | None = File(default=None),
    board_images: list[UploadFile] = File(default=[]),
):
    client_name = client_name.strip()
    sales_rep = sales_rep.strip()
    contract_number = contract_number.strip()

    if ado_file and ado_file.filename and not ado_file.filename.lower().endswith(".pdf"):
        raise HTTPException(status_code=400, detail="Please upload a valid ADO PDF.")

    if not excel_file.filename.lower().endswith((".xlsx", ".xlsm", ".xltx", ".xltm")):
        raise HTTPException(status_code=400, detail="Please upload a valid Excel file.")

    try:
        ado_values = {
            "client_name": "",
            "sales_rep": "",
            "contract_number": "",
        }

        if ado_file and ado_file.filename:
            ado_bytes = await ado_file.read()
            ado_values = extract_ado_contract_fields(ado_bytes)

        final_client_name = choose_contract_value(client_name, ado_values["client_name"], "Client Name")
        final_sales_rep = choose_contract_value(sales_rep, ado_values["sales_rep"], "Sales Representative")
        final_contract_number = choose_contract_value(contract_number, ado_values["contract_number"], "Contract Number")

        file_bytes = await excel_file.read()
        month_year = extract_month_year_from_excel(file_bytes)
        board_rows = extract_board_rows(file_bytes)
        uploaded_images = await read_uploaded_images(board_images)

        pptx_stream = build_pcr_pptx(
            client_name=final_client_name,
            month_year=month_year,
            contract_number=final_contract_number,
            sales_rep=final_sales_rep,
            board_rows=board_rows,
            uploaded_images=uploaded_images,
        )
    except Exception as exc:
        raise HTTPException(status_code=400, detail=str(exc)) from exc

    safe_client = clean_filename_part(final_client_name).replace(" ", "_")
    safe_month = month_year.replace(" ", "_")
    output_filename = f"PCR_Report_{safe_client}_{safe_month}.pptx"

    headers = {
        "Content-Disposition": f'attachment; filename="{output_filename}"'
    }

    return StreamingResponse(
        pptx_stream,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers=headers,
    )


if __name__ == "__main__":
    uvicorn.run("pcr_app:app", host="0.0.0.0", port=8000, reload=True)
