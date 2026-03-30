from io import BytesIO
from pathlib import Path
from datetime import datetime
import re
import copy

from openpyxl import load_workbook
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image, UnidentifiedImageError
from pypdf import PdfReader

from rep_data import REP_DATA


BASE_DIR = Path(__file__).resolve().parent
TEMPLATES_DIR = BASE_DIR / "templates"
PPTX_TEMPLATE_PATH = TEMPLATES_DIR / "PCR - Template.pptx"


def clean_filename_part(value: str) -> str:
    value = value.strip()
    value = re.sub(r"[^A-Za-z0-9 _-]", "", value)
    value = re.sub(r"\s+", " ", value)
    return value or "Client"


def normalize_text(value: str) -> str:
    return re.sub(r"\s+", " ", str(value).strip())


def normalize_match_key(value: str) -> str:
    return re.sub(r"[^a-z0-9]+", "", str(value).lower())


def normalize_excel_label(value) -> str:
    if value is None:
        return ""
    text = str(value).replace("\n", " ")
    text = re.sub(r"\s+", " ", text).strip().upper()
    return text


def _parse_date_string(value: str):
    formats = [
        "%d %B %Y",
        "%d %b %Y",
        "%d/%m/%Y",
        "%d-%m-%Y",
        "%Y-%m-%d",
    ]
    for fmt in formats:
        try:
            return datetime.strptime(value, fmt)
        except ValueError:
            pass
    return None


def format_month_year(value) -> str:
    if isinstance(value, datetime):
        return value.strftime("%B %Y")

    if hasattr(value, "strftime"):
        return value.strftime("%B %Y")

    if isinstance(value, str) and value.strip():
        parsed = _parse_date_string(value.strip())
        if parsed:
            return parsed.strftime("%B %Y")

    raise ValueError("Couldn't format Month Year from the Excel ENDED date.")


def format_day_month_year(value) -> str:
    parsed = None

    if isinstance(value, datetime):
        parsed = value
    elif hasattr(value, "strftime"):
        parsed = value
    elif isinstance(value, str) and value.strip():
        parsed = _parse_date_string(value.strip())

    if not parsed:
        raise ValueError(f"Couldn't format date value: {value}")

    return f"{parsed.day} {parsed.strftime('%b')}, {parsed.year}"


def format_impressions(value) -> str:
    if value is None or str(value).strip() == "":
        return ""

    if isinstance(value, (int, float)):
        return f"{int(round(value)):,}"

    raw = str(value).strip().replace(",", "")
    try:
        return f"{int(float(raw)):,}"
    except ValueError:
        return str(value).strip()


def format_currency(value) -> str:
    if value is None or str(value).strip() == "":
        return ""

    if isinstance(value, (int, float)):
        return f"${value:,.2f}"

    raw = str(value).strip().replace("$", "").replace(",", "")
    try:
        return f"${float(raw):,.2f}"
    except ValueError:
        return str(value).strip()


def format_days(value) -> str:
    if value is None or str(value).strip() == "":
        return ""

    if isinstance(value, (int, float)):
        return f"{int(round(value))} Days"

    raw = str(value).strip()
    return f"{raw} Days"


def get_site_top_line(value) -> str:
    if value is None:
        return ""
    text = str(value).strip()
    if not text:
        return ""
    return text.splitlines()[0].strip()


def extract_site_code_from_left_column(value) -> str:
    if value is None:
        return ""

    text = str(value).strip()
    if not text:
        return ""

    return text.upper()


def extract_month_year_from_excel(file_bytes: bytes) -> str:
    workbook = load_workbook(filename=BytesIO(file_bytes), data_only=True)

    for sheet in workbook.worksheets:
        ended_cell = None

        for row in sheet.iter_rows():
            for cell in row:
                if normalize_excel_label(cell.value) == "ENDED":
                    ended_cell = cell
                    break
            if ended_cell:
                break

        if not ended_cell:
            continue

        date_cell = sheet.cell(row=ended_cell.row + 1, column=ended_cell.column)
        return format_month_year(date_cell.value)

    raise ValueError("Couldn't find an 'ENDED' header in the uploaded Excel file.")


def extract_campaign_insights(file_bytes: bytes) -> dict:
    workbook = load_workbook(filename=BytesIO(file_bytes), data_only=True)

    required_labels = {
        "ELAPSED DAYS:": "Length",
        "CAMPAIGN TOTAL:": "Price",
        "TRAFFIC (CARS):": "Cars",
        "IMPRESSIONS:": "Impressions",
    }

    for sheet in workbook.worksheets:
        found = {}

        for row in sheet.iter_rows():
            for cell in row:
                label = normalize_excel_label(cell.value)
                if label not in required_labels:
                    continue

                value_cell = sheet.cell(row=cell.row, column=cell.column + 1)
                found[required_labels[label]] = value_cell.value

        if len(found) == 4:
            return {
                "Length": format_days(found["Length"]),
                "Price": format_currency(found["Price"]),
                "Cars": format_impressions(found["Cars"]),
                "Impressions": format_impressions(found["Impressions"]),
            }

    raise ValueError(
        "Couldn't find ELAPSED DAYS / CAMPAIGN TOTAL / TRAFFIC (CARS) / IMPRESSIONS in the uploaded Excel file."
    )


def extract_board_rows(file_bytes: bytes):
    workbook = load_workbook(filename=BytesIO(file_bytes), data_only=True)

    required_headers = {
        "SITE",
        "STARTED",
        "ENDED",
        "DAYS",
        "IMPRESSIONS",
    }

    for sheet in workbook.worksheets:
        header_row_idx = None
        header_map = {}

        for row in sheet.iter_rows():
            current_map = {}
            for cell in row:
                header_text = normalize_excel_label(cell.value)
                if header_text in required_headers:
                    current_map[header_text] = cell.column

            if required_headers.issubset(current_map.keys()):
                header_row_idx = row[0].row
                header_map = current_map
                break

        if not header_row_idx:
            continue

        rows = []
        current_row = header_row_idx + 1

        while current_row <= sheet.max_row:
            site_name_value = sheet.cell(current_row, header_map["SITE"]).value
            started_value = sheet.cell(current_row, header_map["STARTED"]).value
            ended_value = sheet.cell(current_row, header_map["ENDED"]).value
            days_value = sheet.cell(current_row, header_map["DAYS"]).value
            impressions_value = sheet.cell(current_row, header_map["IMPRESSIONS"]).value

            site_code_value = None
            if header_map["SITE"] > 1:
                site_code_value = sheet.cell(current_row, header_map["SITE"] - 1).value

            if all(
                value is None or str(value).strip() == ""
                for value in [
                    site_name_value,
                    started_value,
                    ended_value,
                    days_value,
                    impressions_value,
                    site_code_value,
                ]
            ):
                current_row += 1
                continue

            site_name = get_site_top_line(site_name_value)
            site_code = extract_site_code_from_left_column(site_code_value)

            if not site_name:
                current_row += 1
                continue

            site_name_and_code = site_name
            if site_code:
                site_name_and_code = f"{site_name} - {site_code}"

            rows.append(
                {
                    "Site Name and Code": site_name_and_code,
                    "Site Code": site_code,
                    "Start Date": format_day_month_year(started_value),
                    "End Date": format_day_month_year(ended_value),
                    "Run Time": f"{str(days_value).strip()} Days",
                    "Impressions": format_impressions(impressions_value),
                }
            )
            current_row += 1

        if rows:
            return rows

    raise ValueError(
        "Couldn't find a valid SITE / STARTED / ENDED / DAYS / IMPRESSIONS section in the uploaded Excel file."
    )


def extract_ado_contract_fields(pdf_bytes: bytes) -> dict:
    try:
        reader = PdfReader(BytesIO(pdf_bytes))
        text_parts = []
        for page in reader.pages:
            page_text = page.extract_text() or ""
            text_parts.append(page_text)
        full_text = "\n".join(text_parts)
    except Exception as exc:
        raise ValueError("Couldn't read the uploaded ADO PDF.") from exc

    if not full_text.strip():
        raise ValueError("The uploaded ADO PDF did not contain extractable text.")

    def extract_field(label: str) -> str:
        pattern = rf"{re.escape(label)}\s*:\s*(.+)"
        match = re.search(pattern, full_text, re.IGNORECASE)
        if not match:
            return ""
        return normalize_text(match.group(1))

    client_name = extract_field("Advertiser Name")
    sales_rep = extract_field("Account Manager")
    contract_number = extract_field("Contract Number")

    return {
        "client_name": client_name,
        "sales_rep": sales_rep,
        "contract_number": contract_number,
    }


def extract_ado_preview_data(pdf_bytes: bytes) -> dict:
    extracted = extract_ado_contract_fields(pdf_bytes)
    return {
        "client_name": extracted.get("client_name", ""),
        "sales_rep": extracted.get("sales_rep", ""),
        "contract_number": extracted.get("contract_number", ""),
    }


def choose_contract_value(manual_value: str, ado_value: str, label: str) -> str:
    manual_value = normalize_text(manual_value or "")
    ado_value = normalize_text(ado_value or "")

    if manual_value:
        return manual_value
    if ado_value:
        return ado_value

    raise ValueError(f"{label} could not be found. Upload an ADO or enter it manually.")


def replace_text_on_slide(slide, replacements: dict):
    normalized_replacements = {
        normalize_text(key): value for key, value in replacements.items()
    }

    for shape in slide.shapes:
        if not hasattr(shape, "text_frame") or not shape.has_text_frame:
            continue

        current_text = normalize_text(shape.text)
        if current_text not in normalized_replacements:
            continue

        new_text = normalized_replacements[current_text]

        text_frame = shape.text_frame
        if not text_frame.paragraphs:
            text_frame.text = new_text
            continue

        first_paragraph = text_frame.paragraphs[0]
        if first_paragraph.runs:
            first_paragraph.runs[0].text = new_text
            for run in first_paragraph.runs[1:]:
                run.text = ""
        else:
            first_paragraph.text = new_text

        for paragraph in text_frame.paragraphs[1:]:
            for run in paragraph.runs:
                run.text = ""
            if not paragraph.runs:
                paragraph.text = ""


def duplicate_slide_safe(prs, slide_index: int):
    source = prs.slides[slide_index]
    blank_layout = prs.slide_layouts[6]
    new_slide = prs.slides.add_slide(blank_layout)

    for shape in source.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            image_stream = BytesIO(shape.image.blob)
            new_slide.shapes.add_picture(
                image_stream,
                shape.left,
                shape.top,
                width=shape.width,
                height=shape.height,
            )
        else:
            new_el = copy.deepcopy(shape.element)
            new_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")

    return new_slide


def move_slide_to_end(prs, slide_index: int):
    slide_id_list = prs.slides._sldIdLst
    slides = list(slide_id_list)
    slide = slides[slide_index]
    slide_id_list.remove(slide)
    slide_id_list.append(slide)


def remove_slide(prs, slide_index: int):
    slide_id_list = prs.slides._sldIdLst
    slides = list(slide_id_list)
    slide_id_list.remove(slides[slide_index])


def find_board_placeholder_picture(slide, slide_width, slide_height):
    pictures = [shape for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE]
    if not pictures:
        return None

    slide_area = slide_width * slide_height
    filtered = [
        pic for pic in pictures
        if (pic.width * pic.height) < (slide_area * 0.80)
    ]

    candidates = filtered if filtered else pictures
    return max(candidates, key=lambda pic: pic.width * pic.height)


def delete_shape(shape):
    sp = shape.element
    sp.getparent().remove(sp)


def fit_image_within_bounds(image_bytes: bytes, bounds_width: int, bounds_height: int):
    try:
        with Image.open(BytesIO(image_bytes)) as img:
            img_width, img_height = img.size
    except UnidentifiedImageError as exc:
        raise ValueError("One of the uploaded image files could not be read.") from exc

    if img_width <= 0 or img_height <= 0:
        raise ValueError("One of the uploaded image files has invalid dimensions.")

    scale = min(bounds_width / img_width, bounds_height / img_height)
    fitted_width = int(img_width * scale)
    fitted_height = int(img_height * scale)

    return fitted_width, fitted_height


def replace_board_placeholder_image(slide, slide_width, slide_height, image_bytes: bytes):
    placeholder = find_board_placeholder_picture(slide, slide_width, slide_height)
    if not placeholder:
        return

    left = placeholder.left
    top = placeholder.top
    width = placeholder.width
    height = placeholder.height

    fitted_width, fitted_height = fit_image_within_bounds(image_bytes, width, height)

    new_left = left + int((width - fitted_width) / 2)
    new_top = top + int((height - fitted_height) / 2)

    delete_shape(placeholder)
    slide.shapes.add_picture(
        BytesIO(image_bytes),
        new_left,
        new_top,
        width=fitted_width,
        height=fitted_height,
    )


def collect_uploaded_images(board_images):
    allowed_extensions = (".jpg", ".jpeg", ".png")
    collected = []

    for image in board_images:
        if not image.filename:
            continue

        if not image.filename.lower().endswith(allowed_extensions):
            raise ValueError(
                f"Invalid image file '{image.filename}'. Only JPG, JPEG and PNG are allowed."
            )

        collected.append(image)

    return collected


async def read_uploaded_images(board_images):
    uploaded = []
    for image in collect_uploaded_images(board_images):
        image_bytes = await image.read()
        uploaded.append(
            {
                "filename": image.filename,
                "match_key": normalize_match_key(Path(image.filename).stem),
                "bytes": image_bytes,
            }
        )
    return uploaded


def find_matching_image_bytes(site_code: str, uploaded_images: list[dict]):
    code_key = normalize_match_key(site_code)
    if not code_key:
        return None

    for image in uploaded_images:
        if code_key in image["match_key"]:
            return image["bytes"]

    return None


def build_pcr_pptx(
    client_name: str,
    month_year: str,
    contract_number: str,
    sales_rep: str,
    campaign_insights: dict,
    board_rows: list,
    uploaded_images: list[dict],
) -> BytesIO:
    if not PPTX_TEMPLATE_PATH.exists():
        raise FileNotFoundError(f"Template not found at: {PPTX_TEMPLATE_PATH}")

    if sales_rep not in REP_DATA:
        raise ValueError("Selected sales rep wasn't recognised.")

    if not board_rows:
        raise ValueError("No board rows were found in the uploaded Excel file.")

    rep = REP_DATA[sales_rep]
    prs = Presentation(str(PPTX_TEMPLATE_PATH))

    if len(prs.slides) < 4:
        raise ValueError(
            "The PPTX template must contain at least 4 slides: cover, campaign insights, board page, contact page."
        )

    cover_slide_index = 0
    insights_slide_index = 1
    board_template_index = 2

    cover_slide = prs.slides[cover_slide_index]
    replace_text_on_slide(
        cover_slide,
        {
            "Client Name": client_name,
            "Month Year": month_year,
            "Contract Number": contract_number,
        },
    )

    insights_slide = prs.slides[insights_slide_index]
    replace_text_on_slide(
        insights_slide,
        {
            "Length": campaign_insights["Length"],
            "Price": campaign_insights["Price"],
            "Cars": campaign_insights["Cars"],
            "Impressions": campaign_insights["Impressions"],
        },
    )

    for row_data in board_rows:
        new_board_slide = duplicate_slide_safe(prs, board_template_index)
        replace_text_on_slide(new_board_slide, row_data)

        matched_image = find_matching_image_bytes(
            row_data.get("Site Code", ""),
            uploaded_images,
        )
        if matched_image:
            replace_board_placeholder_image(
                new_board_slide,
                prs.slide_width,
                prs.slide_height,
                matched_image,
            )

    remove_slide(prs, board_template_index)

    contact_slide = prs.slides[2]
    rep_contact_line = f'{rep["phone"]} | {rep["email"]}'

    replace_text_on_slide(
        contact_slide,
        {
            "Rep Name!": rep["display_name"],
            "Rep Number | Rep Email": rep_contact_line,
            "Rep Number  |  Rep Email": rep_contact_line,
            "Rep Number   |   Rep Email": rep_contact_line,
        },
    )

    move_slide_to_end(prs, 2)

    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output
